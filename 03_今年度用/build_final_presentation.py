#!/usr/bin/env python3
"""最終発表用 PowerPoint を生成する。"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
MEAS = ROOT / "測定_20260818"
FIG = MEAS / "figures"
THEORY = MEAS / "theory_research" / "figures"
GEO = FIG / "発表スライド_地層断面"
DENOISED = FIG / "地点別_denoised" / "stages" / "02_large_d_cut200" / "theory_16_19"

OUT = MEAS / "最終発表_9班_地下中性子.pptx"
OUT_DL = Path.home() / "Downloads" / "プレゼンテーション1_完成版.pptx"

NAVY = RGBColor(0x1A, 0x3A, 0x5C)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x55, 0x55, 0x55)


def _set_title(slide, text: str, *, subtitle: str | None = None) -> None:
    slide.shapes.title.text = text
    tf = slide.shapes.title.text_frame
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY
    if subtitle and len(slide.placeholders) > 1:
        ph = slide.placeholders[1]
        ph.text = subtitle
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY


def _add_bullets(slide, items: list[str], *, left=0.7, top=1.6, width=8.8, height=5.0, size=18) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(8)


def _add_image(slide, path: Path, *, left=0.5, top=1.5, width=9.0) -> None:
    if not path.exists():
        _add_bullets(slide, [f"（図未配置: {path.name}）"], top=2.5)
        return
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def _title_content(prs: Presentation, title: str, bullets: list[str], *, subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_title(slide, title, subtitle=subtitle)
    _add_bullets(slide, bullets)


def _title_image(prs: Presentation, title: str, img: Path, *, caption: str | None = None, width=9.0, top=1.4) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, title)
    _add_image(slide, img, top=top, width=width)
    if caption:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(8.8), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER


def _section(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(8.0), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(8.0), Inches(1.0))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = ACCENT
        p2.alignment = PP_ALIGN.CENTER


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- 表紙 ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "地下での中性子フラックスの測定"
    slide.placeholders[1].text = (
        "KEKサマーチャレンジ 9班\n"
        "教員：岩瀬 広　岸本 祐二　大山 隆宏\n"
        "学生：川嶋 宥翔(名古屋大学)　鈴江 春樹(明治大学)　高橋 源太(金沢大学)\n"
        "　　　戸田 日々輝(東京大学)　西之原 純輝(宮崎大学)\n"
        "TA：星野 稔　落合 勇稀　小川 拓泰"
    )
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if shape == slide.shapes.title:
                    p.font.size = Pt(32)
                    p.font.bold = True
                    p.font.color.rgb = NAVY

    # --- 目次 ---
    _title_content(
        prs,
        "発表の流れ",
        [
            "1. 背景と目的（昨年度の続き）",
            "2. 予想：深度とともに指数減衰するはず",
            "3. 実験：原理・装置・測定地点",
            "4. 結果：教材理論と大きく乖離",
            "5. 考察：地中での中性子生成（ミュオン起源）",
            "6. 完全理論と神岡地下実験室との整合",
            "7. まとめ",
        ],
    )

    _section(prs, "第1章", "背景と目的")

    _title_content(
        prs,
        "昨年度（9班）の結果",
        [
            "標高を変えて中性子フラックスを測定（管理棟1階 → 白根山 Δh≈2000 m）",
            "熱中性子：λ ≈ 1470 m（Λ ≈ 147 g/cm²）",
            "→ 大気中では深度（標高）に応じて指数関数的に変化",
            "しかし地下・覆土下の測定は未実施",
            "今年度の問い：「深くなるほど e^{-x/λ} で減る」は成り立つか？",
        ],
        subtitle="2025年度演習9班の続き",
    )

    _title_content(
        prs,
        "実験の目的",
        [
            "1. 中性子フラックスの地上からの深度（覆土・遮蔽）依存を測定する",
            "2. 教材の減衰則 φ = A₀·exp(−t_eq/λ_c) が underground で成り立つか検証する",
            "3. 乖離があれば物理要因（ミュオン・環境放射能など）を定量化する",
        ],
    )

    _title_content(
        prs,
        "結果の予想（測定前）",
        [
            "教材・1次元化モデル：φ ∝ exp(−t_eq/λ_c)",
            "λ_c ≈ 39 cm（等価コンクリート換算）",
            "A₀ ≈ 3.07×10⁻³ n/cm²/s（地上基準）",
            "→ 深部（PF→Linac3→KEKB）では桁落ちするはず",
        ],
        subtitle="単純指数減衰を予想",
    )

    _section(prs, "第2章", "原理・装置・較正")

    _title_content(
        prs,
        "中性子線について",
        [
            "一次宇宙線（主に陽子）が大気中の原子核と反応して生成",
            "GeV → MeV → keV → eV → 熱中性子へエネルギー低下",
            "地表付近では宇宙線二次中性子が主要な環境中性子源",
        ],
    )

    _title_image(
        prs,
        "測定原理：³He(n,p)T 反応",
        THEORY / "theory_formula_diagram.png",
        caption="Q = 764 keV。全吸収ピークと壁効果（191 keV 下端）を wall 窓で積分",
        width=8.5,
    )

    _title_content(
        prs,
        "フラックスの定義と較正",
        [
            "Φ = R / (ε·S)　[R: NET 計数率, ε: 検出効率, S: 有効面積]",
            "絶対較正：黒鉛パイル + Am-Be 線源（米内ほか 2002, 保健物理 37(2)）",
            "d1 の εS = 74.2 cm² を基準に D1/D2/d2 へ転送",
            "解析窓（wall）：191–764 keV（³He 反応ピーク＋連続部）",
        ],
    )

    _title_content(
        prs,
        "検出器（4 系統）",
        [
            "D1 / d1：裸管（大径 SN1715 / 小径 SN2162）",
            "D2 / d2：ポリエチレン緩衝付き（MeV 領域感度）",
            "He-3 10 atm，MCA8000D で 512 ch スペクトル取得",
            "昨年と同型の d1 管を継続使用 → 年度間比較が可能",
        ],
    )

    _section(prs, "第3章", "実験：KEK地下トンネルでの測定")

    _title_image(
        prs,
        "測定地点（2026-08-18〜25）",
        GEO / "slide04_sites_map.png",
        caption="全加速器ビーム OFF 期間。地上 → 管理棟 → PF → BT → Linac3 → KEKB → linacIRON",
        width=8.8,
    )

    _title_image(
        prs,
        "等価コンクリート厚 t_eq",
        GEO / "slide05_KEKB_3層地質.png",
        caption="土・鉄を密度換算し t_eq = X/ρ_c（ρ_c=2.3 g/cm³）で横軸を統一",
        width=8.5,
        top=1.3,
    )

    _title_content(
        prs,
        "主要地点と遮蔽（垂直積層）",
        [
            "地上：t_eq = 0 cm（基準 A₀ = 3.07×10⁻³ n/cm²/s）",
            "PF：105 cm コンクリート（2.4 mwe）",
            "Linac3：300 cm コンクリート（6.9 mwe）",
            "KEKB：525 cm 換算（12.1 mwe、3 層地質）",
            "linacIRON：728 cm 換算（土100+コン200+鉄150 cm、1 面開口）",
        ],
    )

    _title_image(
        prs,
        "解析の流れ",
        GEO / "slide03_method.png",
        caption="MCA → wall NET CPS → εS で割って絶対フラックス",
        width=8.5,
    )

    _title_image(
        prs,
        "MCA スペクトル例",
        FIG / "地点別" / "D1_20260819_1530_地上" / "全ch_線形_ch0除く_クリップ.png",
        caption="地上 vs 深部地点でスペクトル形状・計数率が大きく異なる",
        width=4.2,
        top=1.5,
    )
    # 2枚目スペクトル（Linac3）を同スライドに追加
    slide = prs.slides[-1]
    _add_image(
        slide,
        FIG / "地点別" / "D1_20260823_1510_Linac3" / "全ch_線形_ch0除く_クリップ.png",
        left=5.0,
        top=1.5,
        width=4.2,
    )

    _section(prs, "第4章", "結果")

    _title_image(
        prs,
        "結果：旧理論 vs 実測",
        DENOISED / "19_全地点_フラックス_絶対_検出器比較_神岡_連続_誤差棒_片対数.png",
        caption="破線：教材理論 φ=A₀e^{-t/λ_c}。深部で理論より上方にデータ点が分布",
        width=9.0,
        top=1.25,
    )

    _title_content(
        prs,
        "旧理論との乖離（定量）",
        [
            "Linac3（300 cm）：理論比 ~10² 倍「高い」",
            "KEKB（525 cm）：~10⁴ 倍「高い」",
            "linacIRON（728 cm）：~10⁷ 倍「高い」、かつ深部で再上昇",
            "t_eq 190–525 cm 付近で φ ≈ 2.5–4.6×10⁻⁴ に平坦化",
            "→ 単純指数減衰では説明不能",
        ],
        subtitle="深いほど「減りすぎ」ではなく「足されている」",
    )

    _title_image(
        prs,
        "旧理論の残差（地点別）",
        THEORY / "residual_by_site.png",
        caption="log₁₀ RMS 残差：旧理論 2.84 dex",
        width=8.5,
    )

    _section(prs, "第5章", "考察")

    _title_content(
        prs,
        "考察：なぜ理論より「上」にずれるか",
        [
            "単純減衰なら深部ほど下方にずれるはず",
            "実測は深部で平坦化・IRON で再上昇 → 別成分が「足されている」",
            "一次宇宙線中性子は深部でほぼ消失",
            "しかしミュオンは深部まで到達し、物質と反応して中性子を生成",
            "→ 地中での「その場生成」が支配的になる",
        ],
    )

    _title_image(
        prs,
        "ミュオンは深部まで残存",
        THEORY / "cosmic_ray_from_sky_diagram.png",
        caption="PHITS linac_cosmic：覆土 −4.3 m でも μ± は地上の ~80% 残存",
        width=8.0,
    )

    _title_image(
        prs,
        "完全理論モデル",
        THEORY / "component_decomposition.png",
        caption="φ_fast = [F₀e^{-x/Λ_h} + C_μ I_μ(x)] × G_fast + 環境放射能床値",
        width=8.5,
        top=1.3,
    )

    _title_image(
        prs,
        "完全理論 vs 実測",
        THEORY / "theory_vs_meas_all_detectors.png",
        caption="log₁₀ RMS 残差：0.49 dex（旧理論 2.84 dex から大幅改善）",
        width=8.8,
        top=1.2,
    )

    _title_image(
        prs,
        "神岡地下実験室との比較",
        DENOISED / "19_全地点_フラックス_絶対_検出器比較_神岡_連続_誤差棒_片対数.png",
        caption="深部の床値は神岡の環境中性子（~10⁻⁵ n/cm²/s 級）と整合的",
        width=9.0,
        top=1.25,
    )

    _section(prs, "第6章", "まとめ")

    _title_content(
        prs,
        "まとめ",
        [
            "KEK 地下トンネルで初めて深度方向の中性子フラックスを系統測定",
            "教材の単一指数減衰 φ=A₀e^{-t/λ_c} は深部で破綻（最大 10⁷ 倍乖離）",
            "深部の平坦化・再上昇はミュオン起源成分＋環境放射能で説明可能",
            "多成分完全理論で実測と神岡地下値の両方と整合（RMS 0.49 dex）",
            "昨年の「上空方向」に続き、「地下方向」の宇宙線中性子減衰を定量化",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "ご清聴ありがとうございました"
    slide.placeholders[1].text = "質問・バックアップスライドへ"

    # --- バックアップ ---
    _section(prs, "バックアップ", "質問対応用")

    _title_content(
        prs,
        "平均自由行程の求め方",
        [
            "λ = Δx / ln(φ₁/φ₂)",
            "Λ = ΔX / ln(φ₁/φ₂)　[X: 質量厚 g/cm²]",
            "昨年：Δh = 2000 m（管理棟1階 → 白根山）",
            "今年：横軸は等価コンクリート厚 t_eq [cm]",
        ],
    )

    _title_image(
        prs,
        "地層プロファイル（つくば）",
        GEO / "slide01_地層プロファイル概要.png",
        width=8.5,
    )

    _title_image(
        prs,
        "PF・Linac の断面",
        GEO / "slide03_PF_linac_コンクリートのみ.png",
        width=8.5,
    )

    _title_image(
        prs,
        "現場写真",
        GEO / "slide05_field_photos.png",
        width=8.5,
    )

    _title_image(
        prs,
        "PHITS シミュレーション",
        GEO / "slide08_simulation.png",
        width=8.5,
    )

    _title_image(
        prs,
        "12 系統フィット（wall / peak / total）",
        THEORY / "theory_vs_meas_peak_total.png",
        width=8.8,
        top=1.2,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    prs.save(str(OUT_DL))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote: {path}")
    print(f"Wrote: {OUT_DL}")
