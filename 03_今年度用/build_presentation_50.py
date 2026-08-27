#!/usr/bin/env python3
"""GenerateImage スライド 50 枚 + 実測データ図を組み合わせた最終発表 PPTX。"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent
MEAS = ROOT / "測定_20260818"
GEN = MEAS / "figures" / "発表_生成"
THEORY = MEAS / "theory_research" / "figures"
FIG = MEAS / "figures"
DENOISED = FIG / "地点別_denoised" / "stages" / "02_large_d_cut200" / "theory_16_19"

OUT = MEAS / "最終発表_9班_50枚.pptx"
OUT_DL = Path.home() / "Downloads" / "プレゼンテーション1_50枚.pptx"

# 50 slides: generated PNG or real data PNG
SLIDES: list[Path | tuple[Path, Path]] = [
    GEN / "s01_title.png",
    GEN / "s02_agenda.png",
    GEN / "s03_ch1_div.png",
    GEN / "s04_last_year.png",
    GEN / "s05_lambda_air.png",
    GEN / "s06_underground_question.png",
    GEN / "s07_purpose.png",
    GEN / "s08_prediction.png",
    GEN / "s09_ch2_div.png",
    GEN / "s10_air_shower.png",
    GEN / "s11_energy_bands.png",
    GEN / "s12_setup.png",
    GEN / "s13_he3_reaction.png",
    GEN / "s14_mca_peak.png",
    GEN / "s15_wall_effect.png",
    GEN / "s16_flux_formula.png",
    GEN / "s17_calibration.png",
    GEN / "s18_detectors.png",
    GEN / "s19_ch3_div.png",
    GEN / "s20_kek_map.png",
    GEN / "s21_timeline.png",
    GEN / "s22_sites_table.png",
    GEN / "s23_ground_section.png",
    GEN / "s24_pf_section.png",
    GEN / "s25_linac3_section.png",
    GEN / "s26_bt_section.png",
    GEN / "s27_kekb_section.png",
    GEN / "s28_iron_section.png",
    GEN / "s29_teq.png",
    GEN / "s30_analysis.png",
    GEN / "s31_field_photos.png",
    GEN / "s32_ch4_div.png",
    DENOISED / "19_全地点_フラックス_絶対_検出器比較_神岡_連続_誤差棒_片対数.png",
    GEN / "s34_discrepancy.png",
    GEN / "s35_flattening.png",
    GEN / "s36_iron_upturn.png",
    THEORY / "residual_by_site.png",
    (
        FIG / "地点別" / "D1_20260819_1530_地上" / "全ch_線形_ch0除く_クリップ.png",
        FIG / "地点別" / "D1_20260823_1510_Linac3" / "全ch_線形_ch0除く_クリップ.png",
    ),
    GEN / "s39_wall_window.png",
    GEN / "s40_ch5_div.png",
    GEN / "s41_logic_flow.png",
    GEN / "s42_muon_survives.png",
    GEN / "s43_muon_neutron.png",
    GEN / "s44_complete_model.png",
    THEORY / "theory_vs_meas_all_detectors.png",
    THEORY / "component_decomposition.png",
    GEN / "s47_rms_improvement.png",
    GEN / "s48_ch6_div.png",
    GEN / "s49_summary.png",
    GEN / "s50_thanks.png",
]


def _add_full_bleed(slide, path: Path, *, left=0.0, width=10.0) -> None:
    slide.shapes.add_picture(str(path), Inches(left), Inches(0), width=Inches(width), height=Inches(7.5))


def build() -> Path:
    assert len(SLIDES) == 50, len(SLIDES)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for i, item in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        if isinstance(item, tuple):
            left_path, right_path = item
            if not left_path.exists() or not right_path.exists():
                raise FileNotFoundError(f"Slide {i}: {left_path} or {right_path}")
            _add_full_bleed(slide, left_path, left=0.0, width=5.0)
            _add_full_bleed(slide, right_path, left=5.0, width=5.0)
        else:
            if not item.exists():
                raise FileNotFoundError(f"Slide {i}: missing {item}")
            _add_full_bleed(slide, item)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    prs.save(str(OUT_DL))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Slides: 50")
    print(f"Wrote: {path}")
    print(f"Wrote: {OUT_DL}")
