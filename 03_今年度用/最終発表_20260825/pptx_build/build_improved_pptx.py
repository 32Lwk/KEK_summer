#!/usr/bin/env python3
"""9班最終発表原稿の改善版PPTXを新規作成する。

元ファイルのスライド内コメント（改善命令）を反映し、
コメント自体は除去したクリーンな発表用スライドを出力する。
"""

from __future__ import annotations

import copy
import re
import shutil
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

SRC = Path("/Users/yuto/Downloads/9班 最終発表原稿 (1).pptx")
OUT = Path("/Users/yuto/KEK_summer/03_今年度用/最終発表_20260825/9班_最終発表原稿_改善版.pptx")
MEDIA = Path("/Users/yuto/KEK_summer/03_今年度用/最終発表_20260825/pptx_build/media/by_slide")
FIGS = Path("/Users/yuto/KEK_summer/03_今年度用/最終発表_20260825/figures")

NSMAP = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
         "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
         "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"}

# 改善コメントと判定するキーワード（本文の正規コンテンツは除外リストで守る）
COMMENT_MARKERS = [
    "男性衝突", "幅広いエネルギーがある程度", "エネルギー分布だした",
    "今後このように呼ぶ", "前提知識として一枚目", "地表に中性子が降り注いでいる",
    "空から降ってくるけどどれくらいか", "数式だせ", "去年何のため",
    "これまでの9班は何のため", "自然放射線や環境放射線を調べたいよ",
    "もったいないのでフラックス", "-merge", "数式は全部いたりっく",
    "質量欠損(?)", "丁寧な説明は？？", "何で方向が逆なの", "あつい断面積",
    "減速材とかディテクターの話", "ゆっくり(落ち着いて)", "波光が横軸",
    "パルスを数えて縦軸", "型番", "家電粒子", "検出効率が必要なのだ",
    "最後の文を平易に", "何かわからないので改善", "図やイラストでわかりやすく",
    "具体的な数値の方がいい", "放出率はいらないけど", "表でCPSと検出効率",
    "これ急だね", "後つぎと一緒の内容", "ほそくぽい", "補足に入れる", "ほそくにex",
    "色々な深さの地下トンネルがあるからKEKでやったんだよ",
    "予測はここだけで十分", "残差で次のページをなくす", "オーダーが2,3個違う",
    "イベントセレクションが", "唐突感がある", "飛躍がある",
    "自然放射線なら高さで変わるのはおかしくね", "低レイトで測定しているので",
    "Maeってなんやねん", "mに換算しろ", "コンクリでした方がいい",
    "ミューオンが減った時dに基底", "Iμは全部含んでる", "論文のソースにおける意味",
    "１の減少に関しての原因", "EXPACSのフラックス欲しいね",
]


def set_run_font(run, size_pt=None, bold=None, italic=None, color=None, name="游ゴシック"):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    # East Asian font
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", name)
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def clear_paragraphs(tf):
    """テキストフレームを空にする（1段落は残す）。"""
    p = tf.paragraphs[0]
    p.clear()
    for para in list(tf.paragraphs)[1:]:
        p_elem = para._p
        p_elem.getparent().remove(p_elem)
    return tf.paragraphs[0]


def set_textbox_text(shape, lines, size=20, bold=False, color=RGBColor(0x22, 0x22, 0x22), align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    first = clear_paragraphs(tf)
    for i, line in enumerate(lines):
        para = first if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        set_run_font(run, size_pt=size, bold=bold, color=color)


def add_textbox(slide, left, top, width, height, lines, size=20, bold=False,
                color=RGBColor(0x22, 0x22, 0x22), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_textbox_text(box, lines if isinstance(lines, list) else [lines],
                     size=size, bold=bold, color=color, align=align)
    return box


def add_title_bar(slide, title: str, accent=RGBColor(0x1A, 0x56, 0x8A)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.65),
                [title], size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))


def add_footer(slide, page: int):
    add_textbox(slide, Inches(12.2), Inches(7.1), Inches(1.0), Inches(0.3),
                [str(page)], size=12, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.RIGHT)


def add_picture_fit(slide, path: Path, left, top, max_w, max_h):
    if not path.exists():
        return None
    pic = slide.shapes.add_picture(str(path), left, top)
    # fit inside box keeping aspect
    aspect = pic.width / pic.height
    box_aspect = max_w / max_h
    if aspect > box_aspect:
        pic.width = max_w
        pic.height = int(max_w / aspect)
    else:
        pic.height = max_h
        pic.width = int(max_h * aspect)
    pic.left = left
    pic.top = top
    return pic


def shape_full_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs)


def is_comment_shape(shape) -> bool:
    """改善命令テキストボックスかどうか。"""
    if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
        return False
    text = shape_full_text(shape).strip()
    if not text:
        return False
    # 明確なコメントマーカー
    if any(m in text for m in COMMENT_MARKERS):
        return True
    # 口語・命令調の短文（本文プレースホルダーは除外）
    name = shape.name or ""
    if "プレースホルダー" in name or name.startswith("タイトル"):
        return False
    informal = ["のでは", "しろ", "だせ", "すべき", "？？", "‼️", "❣️", "なんや"]
    if any(m in text for m in informal) and len(text) < 400:
        # 本文っぽい定型は除外
        if text.startswith("図") or "フラックス" == text[:5]:
            return False
        return True
    return False


def delete_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def clear_non_placeholders(slide, keep_names=None):
    """プレースホルダー以外の図形を削除（スライド再構成用）。"""
    keep_names = keep_names or set()
    for shape in list(slide.shapes):
        name = shape.name or ""
        if name in keep_names:
            continue
        if getattr(shape, "is_placeholder", False):
            continue
        if "プレースホルダー" in name or name.startswith("タイトル"):
            continue
        try:
            delete_shape(shape)
        except Exception:
            pass


def replace_placeholder_text(slide, placeholder_idx_or_name, lines, size=None, italic_math=False):
    """タイトル/コンテンツプレースホルダーの本文を差し替え。"""
    target = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        name = shape.name or ""
        if placeholder_idx_or_name in name or (
            isinstance(placeholder_idx_or_name, int) and getattr(shape, "is_placeholder", False)
            and shape.placeholder_format.idx == placeholder_idx_or_name
        ):
            target = shape
            break
    if target is None:
        # fallback: first content-like placeholder
        for shape in slide.shapes:
            if shape.has_text_frame and "コンテンツ" in (shape.name or ""):
                target = shape
                break
    if target is None:
        return False

    tf = target.text_frame
    tf.word_wrap = True
    # clear
    first = clear_paragraphs(tf)
    for i, line in enumerate(lines):
        para = first if i == 0 else tf.add_paragraph()
        para.level = 0
        run = para.add_run()
        run.text = line
        is_math = italic_math and bool(re.search(r"[ΦεRSλ=×·/]|exp\(|cm", line))
        set_run_font(run, size_pt=size or 24, italic=is_math)
    return True


def set_title(slide, title: str):
    for shape in slide.shapes:
        if shape.has_text_frame and ("タイトル" in (shape.name or "") or (
            getattr(shape, "is_placeholder", False) and shape.placeholder_format.idx == 0
        )):
            tf = shape.text_frame
            first = clear_paragraphs(tf)
            run = first.add_run()
            run.text = title
            set_run_font(run, size_pt=32, bold=True)
            return True
    return False


def remove_comments(prs: Presentation) -> int:
    removed = 0
    for slide in prs.slides:
        for shape in list(slide.shapes):
            try:
                if is_comment_shape(shape):
                    delete_shape(shape)
                    removed += 1
            except Exception:
                continue
    return removed


def blank_layout(prs):
    # prefer blank
    for layout in prs.slide_layouts:
        if "白紙" in layout.name or "Blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def insert_slide_after(prs, after_index: int, layout):
    """after_index (0-based) の直後に白紙スライドを挿入し、そのスライドを返す。"""
    # python-pptx は末尾追加のみ → 追加後にXMLで並べ替え
    slide = prs.slides.add_slide(layout)
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    new = children[-1]
    sldIdLst.remove(new)
    # insert after after_index
    ref = children[after_index]
    # after removing, find position
    children = list(sldIdLst)
    # find ref again
    idx = None
    for i, c in enumerate(children):
        if c is ref or c.get("id") == ref.get("id"):
            idx = i
            break
    if idx is None:
        sldIdLst.append(new)
    else:
        ref_elem = children[idx]
        ref_elem.addnext(new)
    return slide


def move_slide(prs, old_index: int, new_index: int):
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    el = children[old_index]
    sldIdLst.remove(el)
    children = list(sldIdLst)
    if new_index >= len(children):
        sldIdLst.append(el)
    else:
        children[new_index].addprevious(el)


def delete_slide(prs, index: int):
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    sldId = slides[index]
    # remove relationship
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


# --------------- content improvements ---------------

def improve_slide2(slide):
    """宇宙線中性子の導入 + エネルギー分類の予告。"""
    set_title(slide, "宇宙線由来の中性子線について")
    replace_placeholder_text(slide, "コンテンツ", [
        "一次宇宙線（主に陽子）が大気中の原子核と反応し、大気シャワーを起こす",
        "その二次粒子として、電荷のない中性子が地表へ降り注ぐ",
        "高エネルギー中性子は空気中の原子核との弾性衝突などでエネルギーを失い、",
        "熱エネルギー領域まで減速したものが「熱中性子」と呼ばれる",
        "本発表では、特に「熱中性子」と「MeV中性子」の2帯を測定対象とする",
        "（→ 次スライドでエネルギー帯の呼び方を整理）",
    ], size=20)


def improve_slide3(slide):
    """背景：過去の目的 + 減衰式。"""
    set_title(slide, "背景：これまでの9班と昨年の結果")
    replace_placeholder_text(slide, "コンテンツ", [
        "9班の系列：自然放射線・環境放射線の振る舞いを測り、理解する",
        "2025年度：標高（空気厚）の違いによる宇宙線由来中性子フラックスを測定",
        "結果：大気中でフラックスは概ね指数減衰",
        "　熱中性子：Φ ∝ exp(−6.78×10⁻⁴ · x)　（x：空気厚 [m]）",
        "　MeV中性子：Φ ∝ exp(−9.91×10⁻⁴ · x)",
        "問い：地下でも、深度に対して同様に指数減衰するのか？",
    ], size=20, italic_math=True)


def improve_slide4_purpose_flux(slide):
    """目的とフラックス定義を統合（merge指示）。"""
    clear_non_placeholders(slide)
    set_title(slide, "実験の目的とフラックス Φ")
    replace_placeholder_text(slide, "コンテンツ", [
        "目的：地上からの深度に対する中性子フラックスの変化を調べる",
        "",
        "フラックス Φ = R / (ε · S)　[/cm²/s]",
        "　Φ：1秒あたり 1 cm² を通過する中性子数",
        "　R：計数率（単位時間あたりの検出数）[/s]",
        "　ε：検出効率（入射したうち検出できた割合）",
        "　S：検出面積 [cm²]",
    ], size=22, italic_math=True)


def improve_slide5_energy(slide, prs):
    """旧フラックス詳細スライドを、エネルギー帯の定義スライドに差し替え。"""
    clear_non_placeholders(slide)
    set_title(slide, "中性子のエネルギー帯と今回の観測対象")
    replace_placeholder_text(slide, "コンテンツ", [
        "中性子は運動エネルギーによって呼び方が変わる",
        "本実験の観測対象は次の2帯：",
        "　・熱中性子：〜0.025 eV（熱エネルギー領域）",
        "　・MeV中性子：〜1 MeV 前後（減速材付き検出で感度を持つ帯）",
        "高エネルギー（GeV）帯は³He検出器では直接測れない",
        "→ 以後、スライド中の「熱」「MeV」はこの意味で用いる",
    ], size=20)
    energy_fig = FIGS / "中性子エネルギー帯_観測対象強調.png"
    if energy_fig.exists():
        add_picture_fit(slide, energy_fig, Inches(0.5), Inches(4.55), Inches(12.3), Inches(2.6))


def improve_slide6(slide):
    set_title(slide, "中性子検出の原理：³He(n,p)³H")
    # 説明テキストボックスを追加（既存図を活かす）
    add_textbox(slide, Inches(0.3), Inches(5.5), Inches(12.5), Inches(1.7), [
        "反応：³He + n → ¹H (p) + ³H (t)　＋　Q = 0.764 MeV（質量欠損が運動エネルギーへ）",
        "³Heは熱中性子に対して捕獲断面積が大きく、低エネルギー中性子の検出に適する",
        "検出器内ガスで生じた荷電粒子（p, t）が電離 → パルスとして読み出す",
        "※ 反応後のpとtは反対方向へ放出される（運動量保存）。壁近傍では壁効果でパルスが小さくなる（附記）",
    ], size=15)


def improve_slide7(slide):
    set_title(slide, "実験装置")
    # MCA説明の補足を追加
    add_textbox(slide, Inches(8.3), Inches(6.55), Inches(4.7), Inches(0.7), [
        "MCA：横軸＝チャンネル（エネルギー）、縦軸＝計数",
    ], size=13, color=RGBColor(0x1A, 0x56, 0x8A))


def improve_slide8(slide):
    # 既存の説明テキストを差し替え（写真は残す）
    set_title(slide, "なぜ検出効率が必要か")
    replace_placeholder_text(slide, "コンテンツ", [
        "荷電粒子は電離で痕跡を残しやすいが、",
        "中性子は電荷がなく核反応が起きたときだけ見える",
        "→ 入射したすべてを数えられるわけではない",
        "",
        "放出率が既知の基準線源（²⁴¹Am-Be）で ε を較正",
        "グラファイト減速場で熱中性子束を作り、距離から Φ を得る",
        "平易に言うと：既知の「降り方」と比べて、検出器が何割拾えるかを決める",
    ], size=18, italic_math=True)


def improve_slide9(slide):
    set_title(slide, "検出効率の結果")
    # 左側の数値テキストを簡潔に。右の図は維持。
    replace_placeholder_text(slide, "コンテンツ", [
        "30 cm 位置の熱中性子フラックス",
        "　Φ = 21.33 /cm²/s",
        "",
        "計測率 CPS と面積 S から",
        "　ε = CPS / (Φ · S)",
        "",
        "以後の解析では εS_peak を使用：",
    ], size=18, italic_math=True)

    # 既存テキストと重ならない位置に表
    rows, cols = 5, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(5.0), Inches(5.2), Inches(2.0))
    table = table_shape.table
    data = [
        ("検出器", "対象", "εS_peak [cm²]"),
        ("D1", "熱（裸）", "256.3"),
        ("d1", "熱（裸）", "70.05"),
        ("D2", "MeV（PE）", "178.5"),
        ("d2", "MeV（PE）", "27.25"),
    ]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run_font(run, size_pt=12, bold=(r == 0))


def improve_slide11(slide):
    clear_non_placeholders(slide)
    set_title(slide, "結果の予測（減衰の期待）")
    replace_placeholder_text(slide, "コンテンツ", [
        "昨年の空気厚での結果を踏まえ、地下でも一次の減衰として",
        "　Φ(x) ∝ exp(−x / λ)",
        "を期待する（x：深度、λ：平均自由行程）",
        "",
        "次に測定点と等価コンクリート厚を整理し、",
        "実測がこの単純減衰からどれだけ外れるかを見る",
        "（詳細な予測曲線は結果スライドで一度に示す）",
    ], size=22, italic_math=True)


def improve_slide18(slide):
    """等価コンクリート厚の説明は補足寄りに。"""
    set_title(slide, "補足：等価コンクリート厚への換算")
    add_textbox(slide, Inches(0.4), Inches(6.6), Inches(12), Inches(0.6), [
        "土とコンクリートの密度比で厚みをそろえ、以後の横軸を「等価コンクリート厚」に統一する（詳細は附記）",
    ], size=14, color=RGBColor(0x55, 0x55, 0x55))


def improve_slide19(slide):
    set_title(slide, "測定地点（KEK構内の地下トンネル）")
    add_textbox(slide, Inches(0.3), Inches(6.55), Inches(12.8), Inches(0.7), [
        "KEKには深さの異なる地下トンネルが複数ある → 深度依存を一度のキャンペーンで測れる",
    ], size=16, bold=True, color=RGBColor(0x1A, 0x56, 0x8A))


def improve_slide37(slide):
    set_title(slide, "実験結果と予測の比較")
    add_textbox(slide, Inches(0.3), Inches(6.55), Inches(12.8), Inches(0.7), [
        "単純な指数減衰（文献値ベースの予想）から、フラックスが2〜3桁大きくずれる",
        "→ 地下では「宇宙線中性子の単純減衰」だけでは足りない",
    ], size=15, bold=True, color=RGBColor(0xB0, 0x00, 0x20))


def improve_slide38(slide):
    set_title(slide, "実験結果（まとめプロット）")
    add_textbox(slide, Inches(0.3), Inches(6.7), Inches(12.8), Inches(0.5), [
        "残差・イベントセレクションの詳細は附記。本編ではオーダーの不一致に注目する",
    ], size=14, color=RGBColor(0x55, 0x55, 0x55))


def improve_slide39(slide):
    clear_non_placeholders(slide)
    set_title(slide, "不一致の原因をどう考えるか")
    replace_placeholder_text(slide, "コンテンツ", [
        "観測：深度が増えてもフラックスが予想ほど落ちない（＝「床」がある）",
        "",
        "まず疑うべきは、(A) 深度依存で減る成分 と (B) ほぼ一定の成分 の重ね合わせ",
        "　(A) 宇宙線起源の中性子（深さで減衰）",
        "　(B) 深さにほぼ依らない寄与（例：ミューオンが地下で起こす中性子生成）",
        "",
        "岩石中のU等による自発核分裂なども候補だが、",
        "「深さで変わる／変わらない」に分けて考えると整理しやすい",
        "検出器由来（α汚染など）は、MCAで特徴ピークが支配的でないことから主因としにくい",
    ], size=18)


def improve_slide41(slide):
    set_title(slide, "文献：地下中性子フラックス（換算の注意）")
    add_textbox(slide, Inches(6.2), Inches(2.0), Inches(6.5), Inches(4.5), [
        "文献の横軸はしばしば m.w.e.（メートル水当量）",
        "1 m.w.e. ≒ 水 1 m 分の質量厚さ",
        "　≈ 100 g/cm²",
        "",
        "コンクリート換算の目安：",
        "ρ_c ≈ 2.3 g/cm³ より",
        "1 m.w.e. ≈ 100 / 2.3 ≈ 43 cm コンクリート",
        "",
        "本実験の横軸（等価コンクリート厚）と",
        "見比べるときは、この換算でそろえる",
        "",
        "出典：J. T. Fabryka-Martin (1988)",
    ], size=16)


def improve_slide42(slide):
    set_title(slide, "ミューオン起源中性子の定式化")
    add_textbox(slide, Inches(0.4), Inches(5.8), Inches(12.5), Inches(1.4), [
        "I_μ(d)：深さ d でのミューオン強度。衝突による損失も電離損失による停止も含めた「生き残る強度」",
        "Y_n：1ミューオンあたりの中性子生成収量（平均エネルギーに依存）",
        "N(d) = I_μ(d) × Y_n(Ē_μ) × Δx　で、地下での中性子生成を表す（文献[1][2]）",
    ], size=15)


def improve_closing(slide):
    set_title(slide, "")
    replace_placeholder_text(slide, "コンテンツ", [
        "ご清聴ありがとうございました",
    ], size=40)


def improve_slide53_backup(slide):
    set_title(slide, "バックアップ：エネルギー帯と測定の意味")
    replace_placeholder_text(slide, "コンテンツ", [
        "熱中性子とMeV中性子を分けて測る理由",
        "　・減速・熱化の過程でスペクトルが変わるため、帯ごとの振る舞いが異なる",
        "　・深い地下ではミューオン起因の成分が相対的に効き、帯によって見え方が変わる",
        "",
        "本質的には、高エネルギー中性子が減速した末の成分を³Heで見ている",
        "（EXPACS等の地上スペクトルとの対応は今後の課題）",
    ], size=18)


def build():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(SRC, OUT)
    prs = Presentation(str(OUT))

    removed = remove_comments(prs)
    print(f"Removed comment shapes: {removed}")

    slides = list(prs.slides)

    # index は 0-based。元の番号に対応
    improve_slide2(slides[1])
    improve_slide3(slides[2])
    improve_slide4_purpose_flux(slides[3])
    # slide5 (index4): 旧フラックス詳細 → エネルギー帯定義へ
    improve_slide5_energy(slides[4], prs)
    improve_slide6(slides[5])
    improve_slide7(slides[6])
    improve_slide8(slides[7])
    improve_slide9(slides[8])
    improve_slide11(slides[10])
    improve_slide18(slides[17])
    improve_slide19(slides[18])
    improve_slide37(slides[36])
    improve_slide38(slides[37])
    improve_slide39(slides[38])
    improve_slide41(slides[40])
    improve_slide42(slides[41])
    improve_closing(slides[48])
    if len(slides) > 52:
        improve_slide53_backup(slides[52])

    # スライド5の旧「フラックス概念」図形のうち、明らかに不要なコメントは除去済み。
    # タイトルが "tigau" だった箇所は improve_slide19 で修正済み。

    # 数式イタリック：スライド5旧（今はエネルギー）以外のフラックス式は slide4 で対応。
    # 「ご清聴ありがと」修正済み。

    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
